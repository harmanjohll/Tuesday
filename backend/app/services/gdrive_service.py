"""Google Drive API client.

Uses the same OAuth tokens as Gmail (expanded scopes).
Read files, search, and upload documents.
"""

from __future__ import annotations

import base64
import logging
from datetime import datetime, timedelta, timezone

import httpx

logger = logging.getLogger("tuesday.gdrive")

DRIVE_BASE = "https://www.googleapis.com/drive/v3"
TIMEOUT = 15.0


async def _get_token() -> str | None:
    from app.routers.auth_gmail import load_tokens, refresh_access_token
    tokens = load_tokens()
    if not tokens:
        return None
    refreshed = await refresh_access_token()
    return refreshed or tokens.get("access_token")


def _check() -> str | None:
    from app.routers.auth_gmail import load_tokens
    from app.config import settings
    if not settings.google_client_id:
        return "Google app not configured."
    if not load_tokens():
        return "Google not connected. Visit /auth/gmail to log in (covers Drive too)."
    return None


async def _drive_request(method: str, path: str, **kwargs) -> dict | str:
    from app.routers.auth_gmail import refresh_access_token
    token = await _get_token()
    if not token:
        return "No valid Google token."

    headers = {"Authorization": f"Bearer {token}"}
    async with httpx.AsyncClient(timeout=TIMEOUT) as client:
        resp = await client.request(method, f"{DRIVE_BASE}{path}", headers=headers, **kwargs)

    if resp.status_code == 401:
        new_token = await refresh_access_token()
        if not new_token:
            return "Google auth expired. Re-login at /auth/gmail."
        headers["Authorization"] = f"Bearer {new_token}"
        async with httpx.AsyncClient(timeout=TIMEOUT) as client:
            resp = await client.request(method, f"{DRIVE_BASE}{path}", headers=headers, **kwargs)

    if resp.status_code not in (200, 201):
        return f"Drive API error: {resp.status_code} {resp.text[:200]}"
    return resp.json()


async def list_files(inp: dict) -> str:
    if err := _check():
        return err

    query = inp.get("query", "")
    max_results = min(inp.get("max_results", 15), 25)
    folder_id = inp.get("folder_id", "")

    q_parts = ["trashed = false"]
    if query:
        q_parts.append(f"name contains '{query}'")
    if folder_id:
        q_parts.append(f"'{folder_id}' in parents")

    data = await _drive_request("GET", "/files", params={
        "q": " and ".join(q_parts),
        "pageSize": max_results,
        "fields": "files(id,name,mimeType,modifiedTime,size)",
        "orderBy": "modifiedTime desc",
    })

    if isinstance(data, str):
        return data

    files = data.get("files", [])
    if not files:
        return f"No files found{' matching ' + query if query else ''}."

    lines = [f"Drive files ({len(files)}):"]
    for f in files:
        name = f.get("name", "Untitled")
        mime = f.get("mimeType", "")
        modified = f.get("modifiedTime", "")[:10]
        size = f.get("size", "")
        size_str = f" ({_human_size(int(size))})" if size else ""

        # Simplify mime type
        type_label = mime.split(".")[-1] if "google-apps" in mime else mime.split("/")[-1]
        lines.append(f"- {name} [{type_label}]{size_str} (modified {modified}, id: {f['id']})")

    return "\n".join(lines)


async def read_file(inp: dict) -> str:
    if err := _check():
        return err

    file_id = inp["file_id"]
    token = await _get_token()
    if not token:
        return "No valid Google token."

    headers = {"Authorization": f"Bearer {token}"}

    async with httpx.AsyncClient(timeout=TIMEOUT) as client:
        # First get file metadata
        meta_resp = await client.get(
            f"{DRIVE_BASE}/files/{file_id}",
            headers=headers,
            params={"fields": "name,mimeType,size"},
        )
        if meta_resp.status_code != 200:
            return f"File not found: {meta_resp.status_code}"

        meta = meta_resp.json()
        mime = meta.get("mimeType", "")
        name = meta.get("name", "file")

        # Google Docs/Sheets/Slides → export as plain text or CSV
        if "google-apps.document" in mime:
            resp = await client.get(
                f"{DRIVE_BASE}/files/{file_id}/export",
                headers=headers,
                params={"mimeType": "text/plain"},
            )
        elif "google-apps.spreadsheet" in mime:
            resp = await client.get(
                f"{DRIVE_BASE}/files/{file_id}/export",
                headers=headers,
                params={"mimeType": "text/csv"},
            )
        elif "google-apps.presentation" in mime:
            resp = await client.get(
                f"{DRIVE_BASE}/files/{file_id}/export",
                headers=headers,
                params={"mimeType": "text/plain"},
            )
        else:
            # Regular file — download content
            resp = await client.get(
                f"{DRIVE_BASE}/files/{file_id}",
                headers=headers,
                params={"alt": "media"},
            )

    if resp.status_code != 200:
        return f"Error reading file: {resp.status_code}"

    content = resp.text
    if len(content) > 10000:
        content = content[:10000] + "\n... (truncated)"

    return f"File: {name}\n\n{content}"


async def search_files(inp: dict) -> str:
    """Search Drive files by content or name."""
    if err := _check():
        return err

    query = inp.get("query", "")
    if not query:
        return "No search query provided."

    data = await _drive_request("GET", "/files", params={
        "q": f"fullText contains '{query}' and trashed = false",
        "pageSize": 10,
        "fields": "files(id,name,mimeType,modifiedTime)",
        "orderBy": "modifiedTime desc",
    })

    if isinstance(data, str):
        return data

    files = data.get("files", [])
    if not files:
        return f"No files found containing '{query}'."

    lines = [f"Files matching '{query}' ({len(files)}):"]
    for f in files:
        name = f.get("name", "Untitled")
        modified = f.get("modifiedTime", "")[:10]
        lines.append(f"- {name} (modified {modified}, id: {f['id']})")

    return "\n".join(lines)


async def upload_file(inp: dict) -> str:
    """Upload a file to Google Drive."""
    if err := _check():
        return err

    from pathlib import Path
    from app.config import settings
    import json

    file_path = inp.get("file_path", "")
    filename = inp.get("filename", "")
    folder_id = inp.get("folder_id")
    mime_type = inp.get("mime_type", "application/octet-stream")

    if not file_path:
        return "Error: No file_path provided."

    # Resolve path: relative paths are in outputs dir
    if file_path.startswith("/"):
        local_path = Path(file_path)
    else:
        local_path = Path(settings.outputs_dir) / file_path

    if not local_path.exists():
        return f"Error: File not found at {local_path}"

    if not filename:
        filename = local_path.name

    token = await _get_token()
    if not token:
        return "No valid Google token. Visit /auth/gmail to log in."

    # Build multipart upload body
    metadata = {"name": filename}
    if folder_id:
        metadata["parents"] = [folder_id]

    file_bytes = local_path.read_bytes()
    boundary = "tuesday_upload_boundary"
    meta_json = json.dumps(metadata)

    body = (
        f"--{boundary}\r\n"
        f"Content-Type: application/json; charset=UTF-8\r\n\r\n"
        f"{meta_json}\r\n"
        f"--{boundary}\r\n"
        f"Content-Type: {mime_type}\r\n\r\n"
    ).encode() + file_bytes + f"\r\n--{boundary}--".encode()

    headers = {
        "Authorization": f"Bearer {token}",
        "Content-Type": f"multipart/related; boundary={boundary}",
    }

    async with httpx.AsyncClient(timeout=120.0) as client:
        resp = await client.post(
            "https://www.googleapis.com/upload/drive/v3/files?uploadType=multipart",
            headers=headers,
            content=body,
        )

        # Retry on auth expiry
        if resp.status_code == 401:
            from app.routers.auth_gmail import refresh_access_token
            new_token = await refresh_access_token()
            if not new_token:
                return "Error: Google Drive auth expired. Please re-authenticate at /auth/gmail"
            headers["Authorization"] = f"Bearer {new_token}"
            resp = await client.post(
                "https://www.googleapis.com/upload/drive/v3/files?uploadType=multipart",
                headers=headers,
                content=body,
            )

        if resp.status_code not in (200, 201):
            return f"Error uploading to Drive: {resp.status_code} - {resp.text[:200]}"

        result = resp.json()
        file_id = result.get("id", "unknown")
        file_name = result.get("name", filename)
        logger.info(f"Uploaded {file_name} to Drive (id: {file_id})")
        return f"Uploaded '{file_name}' to Google Drive. File ID: {file_id}."


async def list_folder_contents(folder_name: str) -> list[dict] | str:
    """Find a folder by name and list all files in it."""
    if err := _check():
        return err

    # Find the folder
    data = await _drive_request("GET", "/files", params={
        "q": f"name = '{folder_name}' and mimeType = 'application/vnd.google-apps.folder' and trashed = false",
        "pageSize": 1,
        "fields": "files(id,name)",
    })
    if isinstance(data, str):
        return data

    folders = data.get("files", [])
    if not folders:
        return f"Folder '{folder_name}' not found in Google Drive."

    folder_id = folders[0]["id"]

    # List files in the folder
    data = await _drive_request("GET", "/files", params={
        "q": f"'{folder_id}' in parents and trashed = false",
        "pageSize": 50,
        "fields": "files(id,name,mimeType,modifiedTime,size)",
        "orderBy": "modifiedTime desc",
    })
    if isinstance(data, str):
        return data

    return data.get("files", [])


async def read_file_extended(file_id: str, max_chars: int = 30000) -> str:
    """Read a file with a higher character limit (for analysis use cases)."""
    if err := _check():
        return err

    token = await _get_token()
    if not token:
        return "No valid Google token."

    headers = {"Authorization": f"Bearer {token}"}

    async with httpx.AsyncClient(timeout=60.0) as client:
        meta_resp = await client.get(
            f"{DRIVE_BASE}/files/{file_id}",
            headers=headers,
            params={"fields": "name,mimeType,size"},
        )

        # Handle token refresh on 401
        if meta_resp.status_code == 401:
            from app.routers.auth_gmail import refresh_access_token
            new_token = await refresh_access_token()
            if not new_token:
                return "Google auth expired. Re-login at /auth/gmail."
            headers["Authorization"] = f"Bearer {new_token}"
            meta_resp = await client.get(
                f"{DRIVE_BASE}/files/{file_id}",
                headers=headers,
                params={"fields": "name,mimeType,size"},
            )

        if meta_resp.status_code != 200:
            return f"Error: file not found ({meta_resp.status_code})"

        meta = meta_resp.json()
        mime = meta.get("mimeType", "")
        name = meta.get("name", "file")

        # --- Google-native formats: export as plain text ---
        if "google-apps.document" in mime:
            resp = await client.get(
                f"{DRIVE_BASE}/files/{file_id}/export",
                headers=headers,
                params={"mimeType": "text/plain"},
            )
            if resp.status_code == 200:
                content = resp.text
            else:
                return f"Error: could not export {name} ({resp.status_code})"

        elif "google-apps.spreadsheet" in mime:
            resp = await client.get(
                f"{DRIVE_BASE}/files/{file_id}/export",
                headers=headers,
                params={"mimeType": "text/csv"},
            )
            if resp.status_code == 200:
                content = resp.text
            else:
                return f"Error: could not export {name} ({resp.status_code})"

        elif "google-apps.presentation" in mime:
            resp = await client.get(
                f"{DRIVE_BASE}/files/{file_id}/export",
                headers=headers,
                params={"mimeType": "text/plain"},
            )
            if resp.status_code == 200:
                content = resp.text
            else:
                return f"Error: could not export {name} ({resp.status_code})"

        # --- Uploaded DOCX: download binary + extract text ---
        elif mime in (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/msword",
        ):
            resp = await client.get(
                f"{DRIVE_BASE}/files/{file_id}",
                headers=headers,
                params={"alt": "media"},
            )
            if resp.status_code != 200:
                return f"Error: could not download {name} ({resp.status_code})"
            content = _extract_docx_text(resp.content, name)

        # --- Uploaded PPTX: download binary + extract text ---
        elif mime in (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/vnd.ms-powerpoint",
        ):
            resp = await client.get(
                f"{DRIVE_BASE}/files/{file_id}",
                headers=headers,
                params={"alt": "media"},
            )
            if resp.status_code != 200:
                return f"Error: could not download {name} ({resp.status_code})"
            content = _extract_pptx_text(resp.content, name)

        # --- Uploaded PDF: download binary + extract text ---
        elif mime == "application/pdf":
            resp = await client.get(
                f"{DRIVE_BASE}/files/{file_id}",
                headers=headers,
                params={"alt": "media"},
            )
            if resp.status_code != 200:
                return f"Error: could not download {name} ({resp.status_code})"
            content = _extract_pdf_text(resp.content, name)

        # --- Plain text files: download directly ---
        elif mime.startswith("text/"):
            resp = await client.get(
                f"{DRIVE_BASE}/files/{file_id}",
                headers=headers,
                params={"alt": "media"},
            )
            if resp.status_code != 200:
                return f"Error: could not download {name} ({resp.status_code})"
            content = resp.text

        else:
            return f"Skipped {name}: unsupported file type ({mime})"

    if len(content) > max_chars:
        content = content[:max_chars] + "\n... (truncated)"

    return f"=== {name} ===\n{content}"


def _extract_docx_text(data: bytes, name: str) -> str:
    """Extract text from a DOCX file's binary content."""
    try:
        import io
        from docx import Document
        doc = Document(io.BytesIO(data))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs) if paragraphs else f"(No readable text in {name})"
    except Exception as e:
        logger.error(f"DOCX extraction failed for {name}: {e}")
        return f"(Could not extract text from {name}: {e})"


def _extract_pptx_text(data: bytes, name: str) -> str:
    """Extract text from a PPTX file's binary content."""
    try:
        import io
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            texts.append(para.text.strip())
            if texts:
                slides_text.append(f"[Slide {i}]\n" + "\n".join(texts))
        return "\n\n".join(slides_text) if slides_text else f"(No readable text in {name})"
    except Exception as e:
        logger.error(f"PPTX extraction failed for {name}: {e}")
        return f"(Could not extract text from {name}: {e})"


def _extract_pdf_text(data: bytes, name: str) -> str:
    """Extract text from a PDF file's binary content."""
    try:
        import io
        from PyPDF2 import PdfReader
        reader = PdfReader(io.BytesIO(data))
        pages_text = []
        for i, page in enumerate(reader.pages, 1):
            text = page.extract_text()
            if text and text.strip():
                pages_text.append(f"[Page {i}]\n{text.strip()}")
        return "\n\n".join(pages_text) if pages_text else f"(No readable text in {name} — may be a scanned image PDF)"
    except Exception as e:
        logger.error(f"PDF extraction failed for {name}: {e}")
        return f"(Could not extract text from {name}: {e})"


def _human_size(size_bytes: int) -> str:
    for unit in ["B", "KB", "MB", "GB"]:
        if size_bytes < 1024:
            return f"{size_bytes:.0f} {unit}"
        size_bytes /= 1024
    return f"{size_bytes:.1f} TB"
