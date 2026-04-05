"""Document generation — creates PPTX, DOCX, and PDF files.

Tuesday generates the content via Claude, then these functions
format it into downloadable Office/PDF documents.
"""

from __future__ import annotations

import logging
import uuid
from pathlib import Path

from app.config import settings

logger = logging.getLogger("tuesday.docgen")


def _outputs_dir() -> Path:
    d = settings.outputs_dir
    d.mkdir(parents=True, exist_ok=True)
    return d


def _file_id() -> str:
    return uuid.uuid4().hex[:12]


async def create_presentation(inp: dict) -> str:
    """Generate a PowerPoint presentation, optionally from a template."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    title = inp.get("title", "Presentation")
    slides_data = inp.get("slides", [])
    template_id = inp.get("template_id", "")

    if not slides_data:
        return "No slides provided. Include a 'slides' array with title and content for each slide."

    # Load from template or create blank
    if template_id:
        from app.services.template_service import get_template_path
        template_path = get_template_path(template_id)
        if template_path:
            prs = Presentation(str(template_path))
            logger.info(f"Using template: {template_id}")
        else:
            return f"Template {template_id} not found. Use list_templates to see available templates."
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    # Find best layout indices
    # Templates may have different layout names - try to match intelligently
    title_layout_idx = 0
    content_layout_idx = 1
    for i, layout in enumerate(prs.slide_layouts):
        name_lower = layout.name.lower()
        if "title slide" in name_lower or name_lower == "title":
            title_layout_idx = i
        elif "title and content" in name_lower or "content" in name_lower:
            content_layout_idx = i

    # Title slide
    slide_layout = prs.slide_layouts[title_layout_idx]
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title:
        slide.shapes.title.text = title
    # Try to set subtitle in placeholder 1
    try:
        if slide.placeholders[1]:
            slide.placeholders[1].text = inp.get("subtitle", "")
    except (KeyError, IndexError):
        pass

    # Content slides
    for s in slides_data:
        slide_layout = prs.slide_layouts[content_layout_idx]
        slide = prs.slides.add_slide(slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = s.get("title", "")
        try:
            if slide.placeholders[1]:
                tf = slide.placeholders[1].text_frame
                tf.text = s.get("content", "")
                # Only override font size if not using a template
                if not template_id:
                    for para in tf.paragraphs:
                        for run in para.runs:
                            run.font.size = Pt(18)
        except (KeyError, IndexError):
            pass

    file_id = _file_id()
    filename = f"{file_id}.pptx"
    path = _outputs_dir() / filename
    prs.save(str(path))

    logger.info(f"Created presentation: {filename} ({len(slides_data)} slides)")
    return f"DOWNLOAD:/documents/download/{file_id}|{title}.pptx|Created presentation with {len(slides_data)} slides."


async def create_word_document(inp: dict) -> str:
    """Generate a Word document."""
    from docx import Document
    from docx.shared import Pt

    title = inp.get("title", "Document")
    sections = inp.get("sections", [])
    style = inp.get("style", "formal")

    if not sections:
        return "No sections provided. Include a 'sections' array with heading and body for each section."

    doc = Document()

    # Title
    doc.add_heading(title, level=0)

    # Optional metadata
    if inp.get("author"):
        doc.add_paragraph(f"Author: {inp['author']}")
    if inp.get("date"):
        doc.add_paragraph(f"Date: {inp['date']}")

    # Sections
    for section in sections:
        heading = section.get("heading", "")
        body = section.get("body", "")
        if heading:
            doc.add_heading(heading, level=1)
        if body:
            para = doc.add_paragraph(body)
            for run in para.runs:
                run.font.size = Pt(11)

    file_id = _file_id()
    filename = f"{file_id}.docx"
    path = _outputs_dir() / filename
    doc.save(str(path))

    logger.info(f"Created document: {filename} ({len(sections)} sections)")
    return f"DOWNLOAD:/documents/download/{file_id}|{title}.docx|Created document with {len(sections)} sections."


async def create_pdf_report(inp: dict) -> str:
    """Generate a PDF report."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import cm
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.enums import TA_CENTER

    title = inp.get("title", "Report")
    sections = inp.get("sections", [])

    if not sections:
        return "No sections provided. Include a 'sections' array with heading and body for each section."

    file_id = _file_id()
    filename = f"{file_id}.pdf"
    path = _outputs_dir() / filename

    doc = SimpleDocTemplate(
        str(path),
        pagesize=A4,
        topMargin=2 * cm,
        bottomMargin=2 * cm,
        leftMargin=2.5 * cm,
        rightMargin=2.5 * cm,
    )

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "CustomTitle",
        parent=styles["Title"],
        fontSize=24,
        spaceAfter=30,
        alignment=TA_CENTER,
    )
    heading_style = ParagraphStyle(
        "CustomHeading",
        parent=styles["Heading1"],
        fontSize=16,
        spaceAfter=12,
        spaceBefore=20,
    )
    body_style = ParagraphStyle(
        "CustomBody",
        parent=styles["Normal"],
        fontSize=11,
        spaceAfter=10,
        leading=16,
    )

    elements = []
    elements.append(Paragraph(title, title_style))
    elements.append(Spacer(1, 20))

    for section in sections:
        heading = section.get("heading", "")
        body = section.get("body", "")
        if heading:
            elements.append(Paragraph(heading, heading_style))
        if body:
            # Split by newlines for proper paragraph handling
            for para_text in body.split("\n"):
                if para_text.strip():
                    elements.append(Paragraph(para_text, body_style))

    doc.build(elements)

    logger.info(f"Created PDF report: {filename} ({len(sections)} sections)")
    return f"DOWNLOAD:/documents/download/{file_id}|{title}.pdf|Created PDF report with {len(sections)} sections."
